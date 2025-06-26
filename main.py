import os
import pdfplumber
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.util import Pt, Inches
import pandas as pd
import matplotlib.pyplot as plt

class Application(tk.Frame):
    def __init__(self, master=None):
        super().__init__(master, bg="#F4F4F4")
        self.master = master
        self.master.title("PDF Processor")
        self.master.geometry("800x600")
        self.pack(fill="both", expand=True)
        self.create_main_widgets()

        # Initialize a list to store all invoice data
        self.all_invoices_data = []

    def create_main_widgets(self):
        # Main widgets for PDF processing
        self.select_pdf_label = tk.Label(self, text="Select PDFs", font=("Comic Sans MS", 20, "bold"), fg="black", bg="#E8F0F2")
        self.select_pdf_label.pack(side="top", pady=20)

        self.select_pdf_button = tk.Button(self, text="Select PDFs", font=("Comic Sans MS", 14, "bold"), bg="#4CAF50", fg="white", activebackground="#45A049", activeforeground="white", command=self.select_pdfs)
        self.select_pdf_button.pack(side="top", pady=20)

        self.process_button = tk.Button(self, text="Process PDFs", font=("Comic Sans MS", 14, "bold"), bg="#4CAF50", fg="white", activebackground="#45A049", activeforeground="white", command=self.process_pdfs)
        self.process_button.pack(side="top", pady=20)

        self.status_label = tk.Label(self, text="", font=("Comic Sans MS", 14), fg="black", bg="#F4F4F4")
        self.status_label.pack(side="top", pady=20)

    def select_pdfs(self):
        # Code to select PDFs
        self.pdf_files = filedialog.askopenfilenames(title="Select PDFs", filetypes=[("PDF Files", "*.pdf")])
        self.status_label.config(text=f"Selected {len(self.pdf_files)} PDFs")
        if self.pdf_files:
            self.output_dir = os.path.dirname(self.pdf_files[0])  # Directory of the selected PDFs

    def process_pdfs(self):
        # Code to process PDFs
        if not self.pdf_files:
            self.status_label.config(text="No PDFs selected")
            return

        for pdf_file in self.pdf_files:
            with pdfplumber.open(pdf_file) as pdf:
                page = pdf.pages[0]
                text = page.extract_text()
                self.all_invoices_data.append(text)

        self.generate_csv()
        self.generate_ppt()

    def generate_csv(self):
        # Code to generate a single CSV file
        combined_data = '\n'.join(self.all_invoices_data)
        csv_file_path = os.path.join(self.output_dir, "invoice_data.csv")
        with open(csv_file_path, 'w') as f:
            f.write(combined_data)
        print(f"CSV file saved to {csv_file_path}")
        
        # Display the location of the CSV file
        self.status_label.config(text=f"CSV file saved at: {csv_file_path}")

    def generate_ppt(self):
        # Code to generate PPT files
        ppt = Presentation()

        # Title slide
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])  # Title slide layout
        title = slide.shapes.title
        title.text = "Invoice Processing"

        subtitle = slide.placeholders[1]
        tf = subtitle.text_frame
        tf.text = "Analyzing the Data"

        # Slide 2: Analyzing the Data
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        title.text = "Analyzing the Data"

        tf = slide.placeholders[1].text_frame
        tf.text = "This section analyzes the extracted invoice data."

        # Slide 3: Why this PPT is generated
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        title.text = "Purpose of this Presentation"

        tf = slide.placeholders[1].text_frame
        tf.text = "This presentation is generated to provide an overview of the invoice data and insights."

        # Slide 4: Why we are processing the invoice and data visualization
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        title.text = "Purpose of Invoice Processing and Data Visualization"

        tf = slide.placeholders[1].text_frame
        tf.text = "We process invoices to extract and visualize data to understand sales trends and improve efficiency."

        # Slide 5: Chart with quantity vs part key
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        title.text = "Invoice Quantity vs Part Key"

        # Generate a sample chart (replace with actual data processing)
        part_keys = ['Part A', 'Part B', 'Part C', 'Part D']
        quantities = [10, 5, 8, 2]
        
        plt.bar(part_keys, quantities)
        plt.xlabel('Part Key')
        plt.ylabel('Quantity Sold')
        plt.title('Invoice Quantity vs Part Key')
        chart_path = os.path.join(self.output_dir, "chart.png")
        plt.savefig(chart_path)
        plt.close()

        # Add chart to slide
        img_path = chart_path
        slide.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(8), height=Inches(4))

        # Slide 6: Improvements in Product Selling and Considerations
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        title.text = "Improvements and Considerations"

        tf = slide.placeholders[1].text_frame
        tf.text = "This section provides recommendations for improving product selling based on the analyzed data."

        # Slide 7: Conclusion
        slide = ppt.slides.add_slide(ppt.slide_layouts[1])  # Title and Content layout
        title = slide.shapes.title
        title.text = "Conclusion"

        tf = slide.placeholders[1].text_frame
        tf.text = "Summary of key points and future steps."

        # Save the PPT file
        ppt_file_path = os.path.join(self.output_dir, "invoice_data.pptx")
        ppt.save(ppt_file_path)
        print(f"PPT file saved to {ppt_file_path}")
        self.status_label.config(text=f"PPT file saved at: {ppt_file_path}")

if __name__ == "__main__":
    root = tk.Tk()
    app = Application(master=root)
    app.mainloop()
